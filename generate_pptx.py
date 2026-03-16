"""
ToolTrack Dashboard — 專案簡報 PPTX 產生器
產出: presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ═══════════════════════════════════════════════════════════════
# THEME COLORS (matching dashboard)
# ═══════════════════════════════════════════════════════════════
BG_DEEP    = RGBColor(0x05, 0x0B, 0x12)
BG_CARD    = RGBColor(0x0D, 0x16, 0x20)
BG_ELEV    = RGBColor(0x11, 0x1C, 0x2A)
CYAN       = RGBColor(0x00, 0xD4, 0xFF)
TEAL       = RGBColor(0x00, 0xB8, 0x94)
AMBER      = RGBColor(0xF3, 0x9C, 0x12)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
BLUE       = RGBColor(0x34, 0x98, 0xDB)
PURPLE     = RGBColor(0x9B, 0x59, 0xB6)
WHITE      = RGBColor(0xF0, 0xF6, 0xFC)
SECONDARY  = RGBColor(0xA8, 0xC8, 0xE8)
MUTED      = RGBColor(0x68, 0x89, 0xA8)
BORDER     = RGBColor(0x1A, 0x2D, 0x42)
BORDER_BR  = RGBColor(0x1E, 0x3A, 0x55)
GOLD       = RGBColor(0xFF, 0xD7, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_DEEP):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.rotation = 0
    sf = shape.fill
    if fill_color:
        sf.solid()
        sf.fore_color.rgb = fill_color
    else:
        sf.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sf = shape.fill
    if fill_color:
        sf.solid()
        sf.fore_color.rgb = fill_color
    else:
        sf.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_para(text_frame, text, font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(4)):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet_list(slide, left, top, width, items, font_size=14, color=SECONDARY):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.space_before = Pt(2)
    return txBox


def add_section_head(slide, left, top, text):
    # Accent bar
    add_rect(slide, left, top + Inches(0.02), Inches(0.08), Inches(0.3), fill_color=CYAN)
    add_text_box(slide, left + Inches(0.2), top, Inches(6), Inches(0.4), text, font_size=16, color=CYAN, bold=True)


def add_top_bar(slide):
    """Add the gradient accent line at the top of each slide."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04), fill_color=CYAN)


def add_slide_number(slide, num, total=7):
    add_text_box(slide, Inches(11.8), Inches(0.2), Inches(1.3), Inches(0.3),
                 f"{num:02d} / {total:02d}", font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_slide_title(slide, icon, title, subtitle):
    # Icon box
    shape = add_shape(slide, Inches(0.6), Inches(0.35), Inches(0.55), Inches(0.55),
                      fill_color=BG_CARD, border_color=CYAN, border_width=Pt(2))
    shape.text_frame.paragraphs[0].text = icon
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].font.color.rgb = CYAN
    # Title
    add_text_box(slide, Inches(1.3), Inches(0.3), Inches(8), Inches(0.55),
                 title, font_size=28, color=WHITE, bold=True)
    # Subtitle
    add_text_box(slide, Inches(1.3), Inches(0.85), Inches(8), Inches(0.35),
                 subtitle, font_size=12, color=MUTED)


def add_card(slide, left, top, width, height, fill=BG_ELEV, border=BORDER, accent_color=None):
    """Rounded rectangle card with optional top accent line."""
    shape = add_shape(slide, left, top, width, height, fill_color=fill, border_color=border)
    if accent_color:
        add_rect(slide, left + Inches(0.02), top, width - Inches(0.04), Inches(0.04), fill_color=accent_color)
    return shape


def shape_text(shape, text, font_size=12, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold


def add_arrow_right(slide, left, top, width=Inches(0.6)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    return shape


def add_callout(slide, left, top, width, text, accent_color=CYAN):
    h = Inches(0.7)
    # Left border
    add_rect(slide, left, top, Inches(0.06), h, fill_color=accent_color)
    # Background
    card = add_rect(slide, left + Inches(0.06), top, width - Inches(0.06), h,
                    fill_color=BG_ELEV, border_color=BORDER)
    # Text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), h - Inches(0.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = SECONDARY
    return txBox


# ═══════════════════════════════════════════════════════════════
# SLIDE 0: COVER
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=CYAN)

# Logo
logo = add_shape(slide, Inches(5.9), Inches(1.4), Inches(1.1), Inches(1.1),
                 fill_color=BG_CARD, border_color=CYAN, border_width=Pt(3))
shape_text(logo, "⬡", font_size=40, color=CYAN)

# Title
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "ToolTrack Dashboard", font_size=42, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, Inches(1), Inches(3.65), Inches(11), Inches(0.5),
             "跨站測試工具使用追蹤平台", font_size=18, color=SECONDARY, align=PP_ALIGN.CENTER)

# Meta info
meta_y = Inches(4.6)
meta_items = ["TPE / XM / FQ 三站追蹤", "30 項測試工具管理", "完整 Log 生命週期"]
total_width = len(meta_items) * Inches(3.5)
start_x = (SLIDE_W - total_width) / 2
for i, item in enumerate(meta_items):
    x = start_x + i * Inches(3.5)
    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.1), meta_y + Inches(0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = CYAN
    dot.line.fill.background()
    # Text
    add_text_box(slide, x + Inches(0.35), meta_y, Inches(3), Inches(0.35),
                 item, font_size=13, color=MUTED, align=PP_ALIGN.LEFT)

# Date
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
             "2026 — Project Briefing", font_size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Dashboard 是什麼工具 — 定位 + 數據
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 1)
add_slide_title(slide, "⬡", "Dashboard 是什麼工具", "WHAT IS TOOLTRACK DASHBOARD")

# 工具定位
add_section_head(slide, Inches(0.6), Inches(1.4), "工具定位")
add_bullet_list(slide, Inches(0.6), Inches(1.85), Inches(7), [
    "跨 3 個 Site（TPE / XM / FQ）的測試工具使用追蹤平台",
    "追蹤並管理 30 項測試工具（HW 硬體類 / SW 軟體類）",
    "提供 QA 團隊、Team Lead、管理層進行工具使用率監控與決策",
], font_size=14)

# 數字卡片
stat_data = [("30", "測試工具"), ("3", "跨站追蹤"), ("4", "核心功能")]
for i, (num, label) in enumerate(stat_data):
    x = Inches(8.5) + i * Inches(1.5)
    card = add_card(slide, x, Inches(1.5), Inches(1.35), Inches(1.2), accent_color=CYAN)
    add_text_box(slide, x + Inches(0.1), Inches(1.65), Inches(1.15), Inches(0.6),
                 num, font_size=32, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.2), Inches(1.15), Inches(0.3),
                 label, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# 核心功能架構圖
add_section_head(slide, Inches(0.6), Inches(3.3), "核心功能架構圖")

# Hub
hub = add_shape(slide, Inches(4.2), Inches(3.8), Inches(4.5), Inches(0.6),
                fill_color=BG_CARD, border_color=CYAN, border_width=Pt(2))
shape_text(hub, "⬡  ToolTrack Dashboard", font_size=16, color=CYAN, bold=True)

# Connector line
add_rect(slide, Inches(6.45), Inches(4.4), Inches(0.03), Inches(0.35), fill_color=CYAN)
# Horizontal line
add_rect(slide, Inches(2.3), Inches(4.75), Inches(8.3), Inches(0.03), fill_color=CYAN)

# Feature cards
features = [
    ("📊", "Overview", "年度使用矩陣\n月別使用次數\n跨年度趨勢"),
    ("🏆", "Rankings", "工具排名 Podium\n各站使用排名\nTester 上傳排名"),
    ("📤", "Upload Log", "拖放上傳 Log\n搜尋 / 篩選 / 排序\nPASS / FAIL 標記"),
    ("⚙️", "Tool Status", "CRUD 工具管理\n新增 / 編輯 / 刪除\n啟用 / 停用控制"),
]

for i, (icon, title, desc) in enumerate(features):
    x = Inches(1.3) + i * Inches(2.8)
    # Vertical connector
    add_rect(slide, x + Inches(1.1), Inches(4.75), Inches(0.03), Inches(0.3), fill_color=CYAN)
    # Card
    card = add_card(slide, x, Inches(5.05), Inches(2.4), Inches(1.8), accent_color=CYAN)
    add_text_box(slide, x + Inches(0.1), Inches(5.15), Inches(2.2), Inches(0.35),
                 icon, font_size=24, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(5.5), Inches(2.2), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(5.85), Inches(2.2), Inches(0.9),
                 desc, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: 解決的核心痛點
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 2)
add_slide_title(slide, "🎯", "解決的核心痛點", "PAIN POINTS & SOLUTIONS")

pain_data = [
    ("🚨", "工具使用率不透明", "過去無法得知哪些工具被頻繁使用、\n哪些閒置，資源分配缺乏數據支持", "✔ 年度矩陣 + 月別統計完整可視化", RED),
    ("⚠️", "未使用工具無警示", "閒置工具佔用資源卻無人注意，\n造成維護成本浪費", "✔ 低使用工具紅色脈衝警示 + 標記", AMBER),
    ("🌐", "跨站協作無追蹤", "TPE / XM / FQ 各站工具使用情況\n各自為政，無法統一管理", "✔ 三站統一 Dashboard + 各站排名", TEAL),
]

for i, (icon, title, desc, solve, accent) in enumerate(pain_data):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(1.5)
    card = add_card(slide, x, y, Inches(3.7), Inches(4.5), accent_color=accent)
    # Icon
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.5),
                 icon, font_size=36, align=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(3.1), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.6), Inches(3.1), Inches(1.2),
                 desc, font_size=13, color=MUTED, align=PP_ALIGN.CENTER)
    # Solution
    sol_card = add_shape(slide, x + Inches(0.3), y + Inches(3.2), Inches(3.1), Inches(0.8),
                         fill_color=BG_DEEP, border_color=accent)
    add_text_box(slide, x + Inches(0.5), y + Inches(3.35), Inches(2.7), Inches(0.5),
                 solve, font_size=13, color=accent, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: 當前架構
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 3)
add_slide_title(slide, "🏗", "前後端部署架構 — 當前", "CURRENT DEPLOYMENT ARCHITECTURE")

# Tech badges
techs = [("React 19", CYAN), ("Vite", TEAL), ("GitHub Pages", PURPLE), ("SPA", BLUE)]
for i, (label, color) in enumerate(techs):
    x = Inches(0.6) + i * Inches(1.6)
    badge = add_shape(slide, x, Inches(1.4), Inches(1.4), Inches(0.35),
                      fill_color=BG_ELEV, border_color=color)
    shape_text(badge, label, font_size=11, color=color, bold=True)

# Current architecture flow
nodes = [
    ("💻", "Frontend", "React + Vite SPA"),
    ("📦", "Static Build", "dist/"),
    ("🌐", "GitHub Pages", "靜態託管"),
]

for i, (icon, title, sub) in enumerate(nodes):
    x = Inches(1.5) + i * Inches(3.8)
    y = Inches(2.5)
    is_first = i == 0
    card = add_shape(slide, x, y, Inches(2.5), Inches(1.3),
                     fill_color=BG_CARD, border_color=CYAN if is_first else BORDER_BR,
                     border_width=Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.1), Inches(0.35),
                 icon, font_size=22, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.1), Inches(0.3),
                 title, font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.1), Inches(0.3),
                 sub, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # Arrow between
    if i < len(nodes) - 1:
        ax = x + Inches(2.5)
        add_arrow_right(slide, ax + Inches(0.3), y + Inches(0.45), Inches(0.8))
        label_text = "build" if i == 0 else "deploy"
        add_text_box(slide, ax + Inches(0.3), y + Inches(0.8), Inches(0.8), Inches(0.3),
                     label_text, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Callout
add_callout(slide, Inches(0.6), Inches(4.3), Inches(11.5),
            "ℹ️ 現況：目前為純前端架構，資料寫死在 JSX 中。已可展示完整 UI/UX，下一步需要 Backend 支持動態資料與 Log 上傳。")

# Future label
add_section_head(slide, Inches(0.6), Inches(5.3), "下一步：完整前後端架構")
add_bullet_list(slide, Inches(0.6), Inches(5.8), Inches(11), [
    "Frontend (React) → Backend API Server (Node.js / Python) → Database (PostgreSQL)",
    "加入 Log Parser 引擎、API Key 驗證機制",
    "所有資料動態從 Database 讀取，不再寫死在前端",
], font_size=13)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: 未來完整架構圖
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 4)
add_slide_title(slide, "🔧", "未來完整架構圖", "FUTURE FULL-STACK ARCHITECTURE")

# --- Frontend Tier ---
fe = add_shape(slide, Inches(5.0), Inches(1.5), Inches(3.0), Inches(1.0),
               fill_color=BG_CARD, border_color=CYAN, border_width=Pt(2))
add_text_box(slide, Inches(5.1), Inches(1.55), Inches(2.8), Inches(0.3),
             "💻", font_size=20, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.1), Inches(1.85), Inches(2.8), Inches(0.25),
             "Frontend", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.1), Inches(2.1), Inches(2.8), Inches(0.25),
             "React + Vite", font_size=11, color=CYAN, align=PP_ALIGN.CENTER)

# Arrow down
add_rect(slide, Inches(6.48), Inches(2.5), Inches(0.04), Inches(0.5), fill_color=CYAN)
add_text_box(slide, Inches(7.0), Inches(2.55), Inches(2.5), Inches(0.3),
             "REST API + API Key", font_size=10, color=MUTED)

# --- Backend Tier ---
be_items = [
    ("⚙️", "API Server", "業務邏輯層", "Node.js / Python", TEAL),
    ("🔍", "Log Parser", "Log 解析引擎", "解析欄位資料", TEAL),
    ("🔒", "Auth", "API Key 驗證", "驗證 + 授權", TEAL),
]
for i, (icon, title, sub, tech, clr) in enumerate(be_items):
    x = Inches(2.0) + i * Inches(3.5)
    y = Inches(3.0)
    card = add_shape(slide, x, y, Inches(2.8), Inches(1.2),
                     fill_color=BG_CARD, border_color=clr, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(2.6), Inches(0.3),
                 icon, font_size=18, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.35), Inches(2.6), Inches(0.25),
                 title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.6), Inches(0.2),
                 sub, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(2.6), Inches(0.2),
                 tech, font_size=10, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

# Connector lines to backend
for i in range(3):
    cx = Inches(3.4) + i * Inches(3.5)
    add_rect(slide, cx, Inches(2.85), Inches(0.04), Inches(0.15), fill_color=CYAN)

# Arrow down from backend
add_rect(slide, Inches(6.48), Inches(4.2), Inches(0.04), Inches(0.5), fill_color=TEAL)
add_text_box(slide, Inches(7.0), Inches(4.25), Inches(2.5), Inches(0.3),
             "ORM / Query", font_size=10, color=MUTED)

# --- Database Tier ---
db_items = [
    ("🗄", "Database", "資料持久層", "PostgreSQL / MySQL", AMBER),
    ("💾", "File Storage", "原始 Log 檔", "S3 / Local Disk", AMBER),
]
for i, (icon, title, sub, tech, clr) in enumerate(db_items):
    x = Inches(3.0) + i * Inches(4.0)
    y = Inches(4.7)
    card = add_shape(slide, x, y, Inches(3.2), Inches(1.2),
                     fill_color=BG_CARD, border_color=clr, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.05), Inches(3.0), Inches(0.3),
                 icon, font_size=18, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.35), Inches(3.0), Inches(0.25),
                 title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(3.0), Inches(0.2),
                 sub, font_size=10, color=MUTED, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.85), Inches(3.0), Inches(0.2),
                 tech, font_size=10, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Tier labels
add_text_box(slide, Inches(0.3), Inches(1.7), Inches(1.5), Inches(0.3),
             "FRONTEND", font_size=9, color=MUTED, bold=True)
add_text_box(slide, Inches(0.3), Inches(3.3), Inches(1.5), Inches(0.3),
             "BACKEND", font_size=9, color=MUTED, bold=True)
add_text_box(slide, Inches(0.3), Inches(5.0), Inches(1.5), Inches(0.3),
             "DATABASE", font_size=9, color=MUTED, bold=True)

# Resource needs
add_section_head(slide, Inches(0.6), Inches(6.2), "所需資源")
resources = "應用伺服器 (2 vCPU / 4GB)  ·  PostgreSQL 15+  ·  檔案儲存 100GB+  ·  CI/CD Pipeline  ·  域名 / SSL"
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.3),
             resources, font_size=12, color=SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Log 解析 + API Key
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 5)
add_slide_title(slide, "🔍", "Backend Log 解析 + API Key", "LOG PARSING & API KEY AUTHENTICATION")

# Log parsing flow
add_section_head(slide, Inches(0.6), Inches(1.4), "Log 解析流程")

parse_steps = [
    ("1", "上傳 Log 檔", "使用者透過 Dashboard\n拖放或選擇檔案"),
    ("2", "API Key 驗證", "Backend 驗證請求\n中附帶的 API Key"),
    ("3", "Log 解析", "解析引擎提取\n結構化欄位資料"),
    ("4", "儲存入庫", "結構化資料寫入\nDatabase 供查詢"),
]

for i, (num, title, desc) in enumerate(parse_steps):
    x = Inches(0.6) + i * Inches(3.1)
    y = Inches(1.9)
    card = add_card(slide, x, y, Inches(2.6), Inches(1.5), accent_color=CYAN)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.05), y + Inches(0.15), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYAN
    circle.line.fill.background()
    shape_text(circle, num, font_size=14, color=BG_DEEP, bold=True)
    # Title
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.3), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Desc
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.3), Inches(0.5),
                 desc, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 3:
        add_arrow_right(slide, x + Inches(2.6), y + Inches(0.55), Inches(0.45))

# Parsed fields
add_section_head(slide, Inches(0.6), Inches(3.7), "Log 解析產出欄位")
fields = ["Tool Name", "Site", "Tester", "Result", "Duration", "Timestamp", "Category", "Team"]
for i, f in enumerate(fields):
    x = Inches(0.8) + i * Inches(1.5)
    badge = add_shape(slide, x, Inches(4.15), Inches(1.35), Inches(0.35),
                      fill_color=BG_ELEV, border_color=CYAN)
    shape_text(badge, f, font_size=11, color=CYAN, bold=True)

# API Key flow
add_section_head(slide, Inches(0.6), Inches(4.8), "API Key 驗證機制")

# Flow: Client -> Backend -> Decision -> Results
api_nodes = [
    ("👤", "Client 上傳", "Header 帶入\nX-API-Key"),
    ("🔒", "Backend 驗證", "比對已註冊的\n合法 API Key"),
]
for i, (icon, title, desc) in enumerate(api_nodes):
    x = Inches(0.6) + i * Inches(3.0)
    y = Inches(5.3)
    card = add_card(slide, x, y, Inches(2.5), Inches(1.3))
    add_text_box(slide, x + Inches(0.1), y + Inches(0.1), Inches(2.3), Inches(0.3),
                 icon, font_size=22, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.45), Inches(2.3), Inches(0.25),
                 title, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.7), Inches(2.3), Inches(0.5),
                 desc, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    if i == 0:
        add_arrow_right(slide, Inches(3.1), y + Inches(0.45), Inches(0.45))

# Decision diamond
diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(6.8), Inches(5.45), Inches(1.0), Inches(1.0))
diamond.fill.solid()
diamond.fill.fore_color.rgb = BG_ELEV
diamond.line.color.rgb = AMBER
diamond.line.width = Pt(2)
shape_text(diamond, "驗證?", font_size=11, color=AMBER, bold=True)

add_arrow_right(slide, Inches(6.3), Inches(5.8), Inches(0.45))

# Results
# Pass
pass_card = add_shape(slide, Inches(8.2), Inches(5.35), Inches(4.5), Inches(0.5),
                      fill_color=BG_ELEV, border_color=TEAL)
shape_text(pass_card, "✔ 合法 Key → 允許上傳，解析 Log，寫入 DB", font_size=12, color=TEAL)

add_arrow_right(slide, Inches(7.8), Inches(5.45), Inches(0.35))

# Fail
fail_card = add_shape(slide, Inches(8.2), Inches(6.05), Inches(4.5), Inches(0.5),
                      fill_color=BG_ELEV, border_color=RED)
shape_text(fail_card, "✘ 非法 Key → 拒絕上傳，回傳 401 Unauthorized", font_size=12, color=RED)

add_arrow_right(slide, Inches(7.8), Inches(6.15), Inches(0.35))


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Log 資料生命週期管理
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 6)
add_slide_title(slide, "🔄", "Log 資料生命週期管理", "DATA LIFECYCLE MANAGEMENT")

# Retention callout
add_callout(slide, Inches(0.6), Inches(1.4), Inches(11.5),
            "建議保留期限：2 年 — 涵蓋完整年度對比週期，兼顧分析需求與儲存成本。超過 2 年移入歸檔，滿 3 年後經審核可永久刪除。")

# Lifecycle stages
stages = [
    ("⚡", "Active", "0 ~ 24 個月", "線上即時查詢\n完整搜尋 / 篩選\nDashboard 展示", TEAL),
    ("📦", "Archive", "24 ~ 36 個月", "移入冷儲存\n按需解壓查詢\n壓縮儲存降本", AMBER),
    ("🗑", "Delete", "36 個月後", "經審核永久刪除\n保留統計摘要\n釋放儲存空間", RED),
]

for i, (icon, title, period, desc, clr) in enumerate(stages):
    x = Inches(0.8) + i * Inches(3.8)
    y = Inches(2.5)
    card = add_shape(slide, x, y, Inches(3.3), Inches(2.2),
                     fill_color=BG_ELEV, border_color=clr, border_width=Pt(2))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.9), Inches(0.4),
                 icon, font_size=28, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.9), Inches(0.3),
                 title, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.85), Inches(2.9), Inches(0.25),
                 period, font_size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), Inches(2.9), Inches(0.8),
                 desc, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # Arrow
    if i < 2:
        add_arrow_right(slide, x + Inches(3.3), y + Inches(0.9), Inches(0.45))

# Storage estimation
add_section_head(slide, Inches(0.6), Inches(5.0), "儲存容量估算")
storage_data = [
    ("~2 KB", "每筆 Log 記錄"),
    ("~50 KB", "原始 Log 檔平均"),
    ("~500 MB", "預估年度資料量"),
    ("~1 GB", "2 年保留容量"),
]
for i, (val, label) in enumerate(storage_data):
    x = Inches(0.8) + i * Inches(3.0)
    y = Inches(5.45)
    card = add_card(slide, x, y, Inches(2.7), Inches(1.0), accent_color=CYAN)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.4), Inches(0.4),
                 val, font_size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(2.4), Inches(0.3),
                 label, font_size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Backup & Compliance (brief)
add_section_head(slide, Inches(0.6), Inches(6.6), "備份與合規")
backup_text = "每日增量備份 · 每週完整備份 · 異地備份 · 存取控制 · 審計軌跡 · 資料脫敏 · 刪除雙重確認"
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             backup_text, font_size=12, color=SECONDARY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: 總結與下一步
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_top_bar(slide)
add_slide_number(slide, 7)
add_slide_title(slide, "✨", "總結與下一步", "SUMMARY & NEXT STEPS")

summary_items = [
    ("✅", "已完成", "前端 Dashboard 完成\nOverview / Upload / Tool Status\n三站 30 工具完整展示\nGitHub Pages 部署上線", CYAN),
    ("🚧", "下一階段", "建置 Backend API Server\n實作 Log Parser 引擎\nAPI Key 驗證機制\n資料庫建置與串接", TEAL),
    ("🎯", "預期效益", "即時工具使用率監控\n自動化 Log 解析入庫\n跨站資源調配依據\n資料驅動決策支持", AMBER),
]

for i, (icon, title, desc, clr) in enumerate(summary_items):
    x = Inches(0.8) + i * Inches(4.0)
    y = Inches(1.6)
    card = add_card(slide, x, y, Inches(3.7), Inches(3.8), accent_color=clr)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.3), Inches(0.5),
                 icon, font_size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.3), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.3), Inches(2.0),
                 desc, font_size=14, color=SECONDARY, align=PP_ALIGN.CENTER)

# Discussion callout
add_callout(slide, Inches(0.6), Inches(5.8), Inches(11.5),
            "💬 需要討論：伺服器資源申請  |  資料庫選型確認  |  時程規劃  |  團隊分工")

# Footer
add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
             "ToolTrack Dashboard — Project Briefing — 2026", font_size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
